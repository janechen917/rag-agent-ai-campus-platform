import json
import re
import logging
import random
from rest_framework import viewsets, status
from rest_framework.decorators import api_view, action, permission_classes
from rest_framework.response import Response
from rest_framework.permissions import IsAuthenticated, AllowAny
from rest_framework.parsers import MultiPartParser, FormParser
from django.conf import settings
from django.db.models import Q, Count, Avg, F
from django.utils import timezone

logger = logging.getLogger(__name__)

from .models import (
    AIConversation,
    AIMessage,
    KnowledgeBase,
    CourseRecommendation,
    Quiz,
    QuizQuestion,
    QuizSubmission,
    DebateMatch,
    DebateRound,
    DebateBadge,
)
from .serializers import (
    AIConversationSerializer, AIMessageSerializer, KnowledgeBaseSerializer,
    ChatRequestSerializer, ChatResponseSerializer, CourseRecommendationSerializer,
    QuizSerializer, QuizStudentSerializer, QuizQuestionSerializer, QuizSubmissionSerializer,
    DebateStartSerializer, DebateAttackSerializer, DebateBadgeSerializer
)
from .ai_engine import ai_service
from courses.models import Course, CourseFile


class AIConversationViewSet(viewsets.ModelViewSet):
    """AI对话视图集"""
    serializer_class = AIConversationSerializer
    permission_classes = [IsAuthenticated]
    
    def get_queryset(self):
        return AIConversation.objects.filter(user=self.request.user)
    
    def perform_create(self, serializer):
        serializer.save(user=self.request.user)


@api_view(['POST'])
@permission_classes([IsAuthenticated])
def ai_chat(request):
    """
    AI聊天接口
    
    POST /api/ai/chat/
    {
        "message": "用户消息",
        "conversation_id": 1,  // 可选，对话ID
        "history": []  // 可选，对话历史
    }
    """
    # 验证请求数据
    request_serializer = ChatRequestSerializer(data=request.data)
    if not request_serializer.is_valid():
        return Response(
            request_serializer.errors,
            status=status.HTTP_400_BAD_REQUEST
        )
    
    message = request_serializer.validated_data['message']
    conversation_id = request_serializer.validated_data.get('conversation_id')
    history = request_serializer.validated_data.get('history', [])
    
    # 获取或创建对话
    if conversation_id:
        try:
            conversation = AIConversation.objects.get(
                id=conversation_id,
                user=request.user
            )
        except AIConversation.DoesNotExist:
            conversation = AIConversation.objects.create(
                user=request.user,
                title=message[:50]
            )
    else:
        conversation = AIConversation.objects.create(
            user=request.user,
            title=message[:50]
        )
    
    # 保存用户消息
    user_message = AIMessage.objects.create(
        conversation=conversation,
        role='user',
        content=message
    )
    
    # 构建课程上下文
    course_id = request_serializer.validated_data.get('course_id')
    course_context = None
    if course_id:
        try:
            course = Course.objects.get(id=course_id)
            course_files = CourseFile.objects.filter(course=course).values(
                'file_name', 'file_type', 'description'
            )
            course_context = {
                'course_name': course.title,
                'course_files': list(course_files),
            }
        except Course.DoesNotExist:
            pass
    
    # 获取回答模式
    mode = request_serializer.validated_data.get('mode', 'socratic')
    
    # 调用AI服务
    try:
        ai_response = ai_service.chat(message, history, course_context=course_context, mode=mode)
        
        # 保存AI回复
        ai_message = AIMessage.objects.create(
            conversation=conversation,
            role='assistant',
            content=ai_response
        )
        
        # 基于对话内容推荐课程
        recommended_courses = _get_course_recommendations(message, request.user)
        
        # 构建响应
        response_data = {
            'response': ai_response,
            'conversation_id': conversation.id,
            'recommended_courses': recommended_courses
        }
        
        response_serializer = ChatResponseSerializer(response_data)
        return Response(response_serializer.data)
    
    except Exception as e:
        # 即使出错也要记录日志，但返回友好的响应而不是错误
        print(f"⚠ AI聊天服务异常: {e}")
        
        # 使用备用响应
        fallback_response = "抱歉，AI服务暂时遇到了一些问题。不过我仍然可以为您提供一些基础的学习建议。请告诉我您想了解的技术或遇到的问题，我会尽力帮助您！"
        
        # 保存备用回复
        ai_message = AIMessage.objects.create(
            conversation=conversation,
            role='assistant',
            content=fallback_response
        )
        
        response_data = {
            'response': fallback_response,
            'conversation_id': conversation.id,
            'recommended_courses': []
        }
        
        response_serializer = ChatResponseSerializer(response_data)
        return Response(response_serializer.data)


import base64


@api_view(['POST'])
@permission_classes([IsAuthenticated])
def ai_chat_with_image(request):
    """
    AI聊天接口 - 支持图片和文档文件上传
    
    POST /api/ai/chat-with-file/
    FormData: message, file, history (JSON string)
    支持文件类型: 图片(jpg/png/gif/webp) / PDF / PPT / Word / TXT
    """
    message = request.data.get('message', '')
    uploaded_file = request.FILES.get('file') or request.FILES.get('image')
    history_str = request.data.get('history', '[]')
    conversation_id = request.data.get('conversation_id')

    if not message and not uploaded_file:
        return Response({'error': '请输入消息或上传文件'}, status=status.HTTP_400_BAD_REQUEST)

    try:
        history = json.loads(history_str) if isinstance(history_str, str) else (history_str or [])
    except (json.JSONDecodeError, TypeError):
        history = []

    # 获取或创建对话
    conversation = None
    if conversation_id:
        try:
            conversation = AIConversation.objects.get(
                id=int(conversation_id),
                user=request.user
            )
        except (AIConversation.DoesNotExist, ValueError, TypeError):
            pass
    if not conversation:
        conversation = AIConversation.objects.create(
            user=request.user,
            title=(message[:50] if message else '文件问答')
        )

    # 保存用户消息
    AIMessage.objects.create(
        conversation=conversation,
        role='user',
        content=message or '[文件]'
    )

    # 处理上传文件
    image_base64 = None
    content_type = 'image/jpeg'
    extracted_text = None
    is_image = False

    if uploaded_file:
        content_type = uploaded_file.content_type or ''
        file_name = uploaded_file.name or ''
        ext = ('.' + file_name.rsplit('.', 1)[-1].lower()) if '.' in file_name else ''
        is_image = content_type.startswith('image/')

        if is_image:
            image_data = uploaded_file.read()
            image_base64 = base64.b64encode(image_data).decode('utf-8')
        else:
            # 文档文件：保存为临时文件并提取文本
            import tempfile
            import os
            allowed_exts = ['.pdf', '.ppt', '.pptx', '.doc', '.docx', '.txt']
            if ext not in allowed_exts:
                return Response({'error': '仅支持图片 / PDF / PPT / Word / TXT 文件'}, status=status.HTTP_400_BAD_REQUEST)
            try:
                with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
                    for chunk in uploaded_file.chunks():
                        tmp.write(chunk)
                    tmp_path = tmp.name
                extracted_text = _extract_file_text(tmp_path)
                os.unlink(tmp_path)
            except Exception as e:
                print(f"⚠ 文件解析失败: {e}")
                extracted_text = None

    # 获取回答模式
    mode = request.data.get('mode', 'socratic')
    if mode not in ('socratic', 'direct'):
        mode = 'socratic'
    
    try:
        if is_image and image_base64:
            ai_response = ai_service.chat_with_image(message, image_base64, content_type, history, mode=mode)
        elif extracted_text:
            # 将文档内容作为上下文拼接到消息中
            doc_prompt = f"以下是学生上传的文档内容：\n\n{extracted_text[:8000]}\n\n学生的问题：{message if message else '请分析以上文档内容，并提供学习建议。'}"
            ai_response = ai_service.chat(doc_prompt, history, mode=mode)
        else:
            ai_response = ai_service.chat(message, history, mode=mode)

        AIMessage.objects.create(
            conversation=conversation,
            role='assistant',
            content=ai_response
        )

        return Response({
            'response': ai_response,
            'conversation_id': conversation.id,
        })
    except Exception as e:
        print(f"⚠ AI文件聊天异常: {e}")
        fallback = "抱歉，AI服务暂时遇到了一些问题。请稍后重试。"
        AIMessage.objects.create(
            conversation=conversation,
            role='assistant',
            content=fallback
        )
        return Response({
            'response': fallback,
            'conversation_id': conversation.id,
        })


@api_view(['POST'])
@permission_classes([IsAuthenticated])
def get_recommendations(request):
    """
    获取课程推荐
    
    POST /api/ai/recommendations/
    """
    user = request.user
    
    # 构建用户画像
    user_profile = {
        'user_id': user.id,
        'learning_history': list(
            user.enrollments.values_list('course__category', flat=True)
        ),
        'preferred_categories': list(
            user.enrollments.values_list('course__category', flat=True).distinct()
        )
    }
    
    try:
        # 使用AI服务生成推荐
        recommendations = ai_service.recommend_courses(user_profile)
        
        # 保存推荐记录
        recommendation_record = CourseRecommendation.objects.create(
            user=user,
            recommended_courses=[r['id'] for r in recommendations],
            reason='基于学习历史的AI推荐'
        )
        
        return Response({
            'recommendations': recommendations,
            'record_id': recommendation_record.id
        })
    
    except Exception as e:
        return Response(
            {'error': f'推荐生成失败: {str(e)}'},
            status=status.HTTP_500_INTERNAL_SERVER_ERROR
        )


@api_view(['POST'])
@permission_classes([IsAuthenticated])
def semantic_search(request):
    """
    语义搜索
    
    POST /api/ai/search/
    {
        "query": "搜索关键词",
        "top_k": 5
    }
    """
    query = request.data.get('query', '')
    top_k = request.data.get('top_k', 5)
    
    if not query:
        return Response(
            {'error': '搜索查询不能为空'},
            status=status.HTTP_400_BAD_REQUEST
        )
    
    try:
        results = ai_service.semantic_search(query, top_k)
        return Response({'results': results})
    
    except Exception as e:
        return Response(
            {'error': f'搜索失败: {str(e)}'},
            status=status.HTTP_500_INTERNAL_SERVER_ERROR
        )


def _get_course_recommendations(message: str, user) -> list:
    """
    基于消息内容推荐相关课程（按标题/描述关键词匹配）
    """
    keywords = [
        'python', 'javascript', 'vue', 'react', 'django', 'node',
        'ai', '机器学习', '深度学习', 'java', 'c++', 'html', 'css',
        '数据库', 'sql', '算法', '数据结构', '前端', '后端', '编程',
    ]
    
    message_lower = message.lower()
    matched_keywords = [kw for kw in keywords if kw in message_lower]
    
    if not matched_keywords:
        return []
    
    # 按标题或描述匹配相关课程
    q_filter = Q()
    for kw in matched_keywords:
        q_filter |= Q(title__icontains=kw) | Q(description__icontains=kw)
    
    courses = Course.objects.filter(
        q_filter,
        is_published=True
    ).exclude(
        id__in=user.enrollments.values_list('course_id', flat=True)
    ).order_by('-rating')[:3]
    
    return [
        {
            'id': course.id,
            'title': course.title,
            'description': course.description,
            'level': course.get_level_display(),
        }
        for course in courses
    ]


DEBATE_TOPICS = [
    '记忆知识点比理解原理更重要',
    '课堂作业应该全部使用开卷形式',
    'AI 应该取代传统编程入门课程',
    '考试成绩应当比项目实践占更高权重',
    '团队协作能力比个人技术能力更关键',
    '大学课程应当取消固定教材',
    '算法课不需要数学基础也能学好',
    '所有课程都应该采用翻转课堂',
]

BADGE_RULES = {
    'first_win': {
        'title': '角斗士初胜',
        'description': '首次在 AI 辩论场击败 AI',
        'icon': '🥇',
    },
    'critical_hit': {
        'title': '逻辑暴击',
        'description': '单回合攻击力达到 90+',
        'icon': '⚔️',
    },
    'knowledge_guardian': {
        'title': '知识守卫者',
        'description': '3 回合以上知识得分平均达到 80+',
        'icon': '🛡️',
    },
}


def _safe_json_object(text: str) -> dict:
    """从AI响应中尽量提取JSON对象"""
    if not text:
        return {}
    match = re.search(r'\{.*\}', text, re.DOTALL)
    if not match:
        return {}
    try:
        return json.loads(match.group())
    except Exception:
        return {}


def _extract_keywords(text: str, limit: int = 10) -> list:
    words = re.findall(r'[A-Za-z]{4,}|[\u4e00-\u9fff]{2,}', text or '')
    uniq = []
    for w in words:
        lw = w.lower()
        if lw not in uniq and len(lw) > 1:
            uniq.append(lw)
        if len(uniq) >= limit:
            break
    return uniq


def _build_course_context(course):
    if not course:
        return ''
    files = CourseFile.objects.filter(course=course).values_list('file_name', flat=True)[:8]
    file_text = '、'.join(list(files)) if files else '无'
    safe_title = course.title or '未命名课程'
    safe_description = (course.description or '')[:500]
    return f"课程：{safe_title}\n课程描述：{safe_description}\n课程文件：{file_text}"


def _generate_debate_claim(topic: str, course_context: str = '') -> str:
    def _is_invalid_claim_text(text: str) -> bool:
        if not text:
            return True
        invalid_markers = [
            '你好！我是AI学习导师',
            '当前状态：',
            'AI智能服务暂时运行在基础模式下',
            '请告诉我：',
        ]
        return any(marker in text for marker in invalid_markers)

    def _build_local_claim(seed_topic: str) -> str:
        templates = [
            '我坚持认为“{topic}”被高估了，短期收益看似明显，但长期会降低学习质量与独立思考能力。',
            '我的立场是：在大多数课堂场景里，“{topic}”并非最优解，因为它牺牲了可迁移能力与知识深度。',
            '我主张“{topic}”应被谨慎限制使用，效率提升只是表象，真正代价是理解力和推理能力被削弱。',
            '我认为“{topic}”不是学习加速器，而是认知捷径陷阱，越依赖越难形成稳定的知识结构。',
        ]
        return random.choice(templates).format(topic=seed_topic)

    prompt = (
        '你是课堂辩论赛中的 AI 对手。请针对下面辩题，给出一句具有争议性的立场陈述。'
        '要求：1）观点要鲜明可反驳；2）字数 40-90；3）不输出编号或解释。\n\n'
        f'辩题：{topic}\n\n'
        f'课程上下文（可选）：\n{course_context[:1200]}'
    )
    try:
        result = ai_service.chat(prompt, [], mode='direct').strip()
        result = re.sub(r'^\s*[\-\d\.、:：]+', '', result)
        if _is_invalid_claim_text(result):
            return _build_local_claim(topic)
        return result[:300] if result else _build_local_claim(topic)
    except Exception:
        return _build_local_claim(topic)


def _evaluate_argument(argument: str, topic: str, ai_claim: str, course_context: str = '') -> dict:
    prompt = f"""你是严谨的辩论裁判。请评估学生反驳质量并给出JSON结果。

辩题：{topic}
AI观点：{ai_claim}
学生反驳：{argument}
课程上下文：{course_context[:1200]}

评分维度（0-100整数）：
- logic_score：论证逻辑
- evidence_score：证据与例子
- knowledge_score：课程知识使用
- structure_score：结构清晰度

攻击力 attack_power 按以下公式计算后取整：
attack_power = logic*0.35 + evidence*0.2 + knowledge*0.3 + structure*0.15

verdict 用一句简短中文点评。

只返回JSON对象：
{{
  "logic_score": 0,
  "evidence_score": 0,
  "knowledge_score": 0,
  "structure_score": 0,
  "attack_power": 0,
  "verdict": ""
}}"""
    try:
        raw = ai_service.chat(prompt, [], mode='direct')
        data = _safe_json_object(raw)
        if data:
            logic = int(max(0, min(100, data.get('logic_score', 0))))
            evidence = int(max(0, min(100, data.get('evidence_score', 0))))
            knowledge = int(max(0, min(100, data.get('knowledge_score', 0))))
            structure = int(max(0, min(100, data.get('structure_score', 0))))
            attack = int(round(logic * 0.35 + evidence * 0.2 + knowledge * 0.3 + structure * 0.15))
            return {
                'logic_score': logic,
                'evidence_score': evidence,
                'knowledge_score': knowledge,
                'structure_score': structure,
                'attack_power': max(0, min(100, int(data.get('attack_power', attack) or attack))),
                'verdict': str(data.get('verdict', '论证具备一定说服力，可继续强化证据。'))[:120],
            }
    except Exception:
        pass

    text = argument or ''
    length = len(text)
    connectors = ['因为', '所以', '因此', '但是', '然而', '首先', '其次', '最后', '例如', '比如', 'if', 'then', 'therefore']
    connector_hits = sum(1 for c in connectors if c.lower() in text.lower())
    logic_score = min(95, 45 + connector_hits * 7 + (12 if length > 180 else 0) + (10 if length > 320 else 0))

    evidence_markers = ['数据', '研究', '实验', '案例', '统计', '论文', '教材', '章节', '图', '表']
    evidence_hits = sum(1 for m in evidence_markers if m in text)
    evidence_score = min(95, 35 + evidence_hits * 10 + (8 if '例如' in text or '比如' in text else 0))

    course_keywords = _extract_keywords(course_context, limit=12)
    kw_hits = sum(1 for k in course_keywords if k and k in text.lower())
    knowledge_score = min(95, 40 + kw_hits * 12)

    paragraphs = [p for p in re.split(r'[\n。！？!?]', text) if p.strip()]
    structure_score = min(95, 40 + min(5, len(paragraphs)) * 8 + (8 if any(w in text for w in ['首先', '其次', '最后']) else 0))

    attack_power = int(round(logic_score * 0.35 + evidence_score * 0.2 + knowledge_score * 0.3 + structure_score * 0.15))
    verdict = '反驳有基础火力，建议加入更具体教材证据，形成闭环论证。'
    if attack_power >= 85:
        verdict = '反驳打击面精准，逻辑与知识点结合非常有力。'
    elif attack_power >= 70:
        verdict = '反驳较有说服力，再补充可验证证据会更强。'

    return {
        'logic_score': int(logic_score),
        'evidence_score': int(evidence_score),
        'knowledge_score': int(knowledge_score),
        'structure_score': int(structure_score),
        'attack_power': int(max(0, min(100, attack_power))),
        'verdict': verdict,
    }


def _generate_ai_counter(topic: str, ai_claim: str, argument: str, score_data: dict) -> str:
    prompt = f"""你是辩论对手。请针对学生反驳给出一句到两句反击。
要求：
1) 语气有竞技感，但不攻击人格；
2) 点名学生论证中的一个薄弱点；
3) 结尾抛出一个追问。

辩题：{topic}
AI观点：{ai_claim}
学生反驳：{argument}
评分：{json.dumps(score_data, ensure_ascii=False)}
"""
    try:
        result = ai_service.chat(prompt, [], mode='direct').strip()
        invalid_markers = ['你好！我是AI学习导师', '当前状态：', '请告诉我：']
        if not result or any(marker in result for marker in invalid_markers):
            return '你的反驳有亮点，但证据还不够硬。你能给出更直接的教材依据吗？'
        return result[:320]
    except Exception:
        return '你的反驳有亮点，但证据还不够硬。你能给出更直接的教材依据吗？'


def _grant_debate_badges(user, match: DebateMatch, latest_round: DebateRound):
    earned_badges = []

    win_count = DebateMatch.objects.filter(student=user, status='won').count()
    if win_count == 1:
        badge, created = DebateBadge.objects.get_or_create(
            student=user,
            code='first_win',
            defaults=BADGE_RULES['first_win'],
        )
        if created:
            earned_badges.append(badge)

    if latest_round.attack_power >= 90:
        badge, created = DebateBadge.objects.get_or_create(
            student=user,
            code='critical_hit',
            defaults=BADGE_RULES['critical_hit'],
        )
        if created:
            earned_badges.append(badge)

    rounds = match.rounds.all()
    if rounds.count() >= 3:
        avg_knowledge = rounds.aggregate(v=Avg('knowledge_score'))['v'] or 0
        if avg_knowledge >= 80:
            badge, created = DebateBadge.objects.get_or_create(
                student=user,
                code='knowledge_guardian',
                defaults=BADGE_RULES['knowledge_guardian'],
            )
            if created:
                earned_badges.append(badge)

    return earned_badges


class KnowledgeBaseViewSet(viewsets.ModelViewSet):
    """知识库视图集"""
    queryset = KnowledgeBase.objects.all()
    serializer_class = KnowledgeBaseSerializer
    permission_classes = [IsAuthenticated]
    
    @action(detail=False, methods=['post'])
    def batch_add(self, request):
        """批量添加知识"""
        items = request.data.get('items', [])
        
        created_items = []
        for item in items:
            kb = KnowledgeBase.objects.create(**item)
            created_items.append(kb)
        
        # 添加到向量数据库
        texts = [item.content for item in created_items]
        metadatas = [{'id': item.id, 'title': item.title} for item in created_items]
        ai_service.add_to_knowledge_base(texts, metadatas)
        
        serializer = self.get_serializer(created_items, many=True)
        return Response(serializer.data, status=status.HTTP_201_CREATED)


def _extract_file_text(file_path):
    """根据文件类型提取文本内容"""
    ext = ('.' + file_path.rsplit('.', 1)[-1].lower()) if '.' in file_path else ''
    if ext in ('.ppt', '.pptx'):
        return _extract_ppt_text(file_path)
    elif ext == '.pdf':
        return _extract_pdf_text(file_path)
    elif ext in ('.doc', '.docx'):
        return _extract_word_text(file_path)
    elif ext == '.txt':
        return _extract_txt_text(file_path)
    else:
        raise ValueError(f'不支持的文件类型: {ext}')


def _extract_ppt_text(file_path):
    """从PPT文件中提取文本内容"""
    from pptx import Presentation
    prs = Presentation(file_path)
    text_content = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_text.append(text)
        if slide_text:
            text_content.append(f"第{slide_num}页:\n" + "\n".join(slide_text))
    return "\n\n".join(text_content)


def _extract_pdf_text(file_path):
    """从PDF文件中提取文本内容"""
    from PyPDF2 import PdfReader
    reader = PdfReader(file_path)
    text_content = []
    for page_num, page in enumerate(reader.pages, 1):
        text = page.extract_text()
        if text and text.strip():
            text_content.append(f"第{page_num}页:\n{text.strip()}")
    return "\n\n".join(text_content)


def _extract_word_text(file_path):
    """从Word文件中提取文本内容"""
    from docx import Document
    doc = Document(file_path)
    text_content = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            text_content.append(text)
    return "\n\n".join(text_content)


def _extract_txt_text(file_path):
    """从TXT文件中提取文本内容"""
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read()


def _generate_quiz_with_ai(ppt_text, question_count):
    """使用AI根据文件内容生成Quiz题目"""
    prompt = f"""请根据以下课件内容，生成{question_count}道四选一的选择题。

课件内容:
{ppt_text[:4000]}

请严格按照以下JSON格式返回，不要包含其他文字：
[
  {{
    "question_text": "题目内容",
    "option_a": "选项A内容",
    "option_b": "选项B内容",
    "option_c": "选项C内容",
    "option_d": "选项D内容",
    "correct_answer": "B",
    "explanation": "解析说明"
  }}
]

要求:
1. 题目要紧密结合课件的核心内容
2. 选项要有迷惑性，但正确答案必须明确
3. 每题都要有简短的解析
4. correct_answer只能是A/B/C/D之一
5. 生成恰好{question_count}道题
6. 非常重要：正确答案必须均匀分布在A、B、C、D四个选项中，不要让所有答案都是同一个选项。请将正确内容随机放在不同选项位置"""

    try:
        ai_response = ai_service.chat(prompt, [], mode='direct')
        # 从响应中提取JSON
        json_match = re.search(r'\[.*\]', ai_response, re.DOTALL)
        if json_match:
            questions = json.loads(json_match.group())
            return questions[:question_count]
    except Exception as e:
        print(f"AI生成Quiz失败: {e}")

    # 备用：如果AI不可用，生成基于PPT内容的简单题目
    return _generate_fallback_questions(ppt_text, question_count)


def _generate_fallback_questions(ppt_text, question_count):
    """在AI不可用时根据PPT文本生成基本题目"""
    import random
    option_labels = ['A', 'B', 'C', 'D']
    lines = [line.strip() for line in ppt_text.split('\n') if line.strip() and not line.startswith('第') and len(line.strip()) > 5]
    questions = []
    for i in range(min(question_count, len(lines))):
        content = lines[i] if i < len(lines) else f"知识点{i+1}"
        options = [
            f"{content[:40]}",
            f"与{content[:20]}无关的内容",
            f"{content[:20]}的反义描述",
            f"以上都不正确",
        ]
        correct_idx = i % 4
        # Rotate options so correct answer lands on different positions
        rotated = options[-correct_idx:] + options[:-correct_idx] if correct_idx else options
        questions.append({
            "question_text": f"以下关于「{content[:30]}」的说法，哪个是正确的？",
            "option_a": rotated[0],
            "option_b": rotated[1],
            "option_c": rotated[2],
            "option_d": rotated[3],
            "correct_answer": option_labels[correct_idx],
            "explanation": f"根据课件内容，正确答案为{option_labels[correct_idx]}。"
        })
    # 补齐不足的题目
    while len(questions) < question_count:
        idx = len(questions)
        correct_idx = idx % 4
        options = [
            "课件中提到的正确描述",
            "不正确的描述",
            "无关的描述",
            "以上都不正确",
        ]
        rotated = options[-correct_idx:] + options[:-correct_idx] if correct_idx else options
        questions.append({
            "question_text": f"关于本课件的第{idx+1}个知识点，以下哪个说法正确？",
            "option_a": rotated[0],
            "option_b": rotated[1],
            "option_c": rotated[2],
            "option_d": rotated[3],
            "correct_answer": option_labels[correct_idx],
            "explanation": "请参考课件内容。"
        })
    return questions[:question_count]


@api_view(['POST'])
@permission_classes([IsAuthenticated])
def upload_and_generate_quiz(request):
    """
    上传文件并生成Quiz

    POST /api/ai/quiz/generate/
    FormData:
        file: PPT/PDF/Word/TXT文件
        title: Quiz标题
        question_count: 题目数量(默认5)
        end_time: 截止时间(ISO格式，可选)
        course_id: 关联课程ID(可选)
    """
    file = request.FILES.get('file')
    if not file:
        return Response({'error': '请上传文件'}, status=status.HTTP_400_BAD_REQUEST)

    # 验证文件类型
    allowed_extensions = ['.ppt', '.pptx', '.pdf', '.doc', '.docx', '.txt']
    file_ext = '.' + file.name.rsplit('.', 1)[-1].lower() if '.' in file.name else ''
    if file_ext not in allowed_extensions:
        return Response({'error': '仅支持 PPT/PPTX/PDF/DOC/DOCX/TXT 文件格式'}, status=status.HTTP_400_BAD_REQUEST)

    # 验证文件大小 (最大100MB)
    if file.size > 100 * 1024 * 1024:
        return Response({'error': '文件大小不能超过100MB'}, status=status.HTTP_400_BAD_REQUEST)

    title = request.data.get('title', file.name.rsplit('.', 1)[0])
    question_count = int(request.data.get('question_count', 5))
    question_count = max(1, min(question_count, 30))  # 限制1-30题
    max_attempts = int(request.data.get('max_attempts', 1))
    max_attempts = max(1, min(max_attempts, 99))  # 限制1-99次
    end_time = request.data.get('end_time') or None
    course_id = request.data.get('course_id') or None

    # 验证课程权限
    course = None
    if course_id:
        try:
            course = Course.objects.get(id=course_id, instructor=request.user)
        except Course.DoesNotExist:
            return Response({'error': '课程不存在或无权操作'}, status=status.HTTP_403_FORBIDDEN)

    # 创建Quiz记录
    try:
        quiz = Quiz.objects.create(
            title=title,
            course=course,
            creator=request.user,
            source_file=file,
            source_file_name=file.name,
            question_count=question_count,
            max_attempts=max_attempts,
            end_time=end_time,
        )
    except Exception as e:
        return Response({'error': f'创建Quiz失败: {str(e)}'}, status=status.HTTP_400_BAD_REQUEST)

    # 提取文件文本
    try:
        ppt_text = _extract_file_text(quiz.source_file.path)
    except Exception as e:
        quiz.delete()
        return Response({'error': f'文件解析失败: {str(e)}'}, status=status.HTTP_400_BAD_REQUEST)

    if not ppt_text.strip():
        quiz.delete()
        return Response({'error': '文件中未找到文本内容'}, status=status.HTTP_400_BAD_REQUEST)

    # 使用AI生成题目
    questions_data = _generate_quiz_with_ai(ppt_text, question_count)

    # 保存题目
    for idx, q in enumerate(questions_data):
        QuizQuestion.objects.create(
            quiz=quiz,
            question_text=q.get('question_text', ''),
            option_a=q.get('option_a', ''),
            option_b=q.get('option_b', ''),
            option_c=q.get('option_c', ''),
            option_d=q.get('option_d', ''),
            correct_answer=q.get('correct_answer', 'A'),
            explanation=q.get('explanation', ''),
            order=idx + 1,
        )

    serializer = QuizSerializer(quiz)
    return Response(serializer.data, status=status.HTTP_201_CREATED)


@api_view(['GET'])
@permission_classes([IsAuthenticated])
def my_quizzes(request):
    """获取教师创建的所有Quiz"""
    quizzes = Quiz.objects.filter(creator=request.user)
    serializer = QuizSerializer(quizzes, many=True)
    return Response(serializer.data)


@api_view(['POST'])
@permission_classes([IsAuthenticated])
def publish_quiz(request, quiz_id):
    """发布Quiz到课程"""
    try:
        quiz = Quiz.objects.get(id=quiz_id, creator=request.user)
    except Quiz.DoesNotExist:
        return Response({'error': 'Quiz不存在或无权操作'}, status=status.HTTP_404_NOT_FOUND)

    course_id = request.data.get('course_id')
    if course_id:
        try:
            course = Course.objects.get(id=course_id, instructor=request.user)
            quiz.course = course
        except Course.DoesNotExist:
            return Response({'error': '课程不存在或无权操作'}, status=status.HTTP_403_FORBIDDEN)

    if quiz.end_time and timezone.now() > quiz.end_time:
        return Response({'error': 'Quiz已截止，无法发布'}, status=status.HTTP_400_BAD_REQUEST)

    quiz.is_published = True
    quiz.save()

    serializer = QuizSerializer(quiz)
    return Response(serializer.data)


@api_view(['GET'])
@permission_classes([IsAuthenticated])
def quiz_detail(request, quiz_id):
    """获取Quiz详情（教师端包含答案，学生端不含）"""
    try:
        quiz = Quiz.objects.get(id=quiz_id)
    except Quiz.DoesNotExist:
        return Response({'error': 'Quiz不存在'}, status=status.HTTP_404_NOT_FOUND)

    # 检查是否是创建者
    if quiz.creator == request.user:
        serializer = QuizSerializer(quiz)
    else:
        # 学生：检查Quiz是否已发布且未过期
        if not quiz.is_published:
            return Response({'error': 'Quiz尚未发布'}, status=status.HTTP_403_FORBIDDEN)
        if quiz.end_time and timezone.now() > quiz.end_time:
            return Response({'error': 'Quiz已截止'}, status=status.HTTP_403_FORBIDDEN)
        serializer = QuizStudentSerializer(quiz, context={'request': request})

    return Response(serializer.data)


@api_view(['GET'])
def quiz_by_share_code(request, share_code):
    """通过分享码获取Quiz（学生端）"""
    try:
        quiz = Quiz.objects.get(share_code=share_code, is_published=True)
    except Quiz.DoesNotExist:
        return Response({'error': 'Quiz不存在或未发布'}, status=status.HTTP_404_NOT_FOUND)

    if quiz.end_time and timezone.now() > quiz.end_time:
        return Response({'error': 'Quiz已截止'}, status=status.HTTP_403_FORBIDDEN)

    if request.user.is_authenticated:
        serializer = QuizStudentSerializer(quiz, context={'request': request})
    else:
        serializer = QuizStudentSerializer(quiz)
    return Response(serializer.data)


@api_view(['POST'])
@permission_classes([IsAuthenticated])
def submit_quiz(request, quiz_id):
    """学生提交Quiz答案"""
    try:
        quiz = Quiz.objects.get(id=quiz_id, is_published=True)
    except Quiz.DoesNotExist:
        return Response({'error': 'Quiz不存在或未发布'}, status=status.HTTP_404_NOT_FOUND)

    if quiz.end_time and timezone.now() > quiz.end_time:
        return Response({'error': 'Quiz已截止'}, status=status.HTTP_403_FORBIDDEN)

    # 检查答题次数限制
    attempts_used = QuizSubmission.objects.filter(quiz=quiz, student=request.user).count()
    if attempts_used >= quiz.max_attempts:
        return Response({'error': f'您已达到最大答题次数({quiz.max_attempts}次)'}, status=status.HTTP_400_BAD_REQUEST)

    answers = request.data.get('answers', {})
    if not answers:
        return Response({'error': '请提供答案'}, status=status.HTTP_400_BAD_REQUEST)

    # 计算分数
    questions = quiz.questions.all()
    total = questions.count()
    correct = 0
    for question in questions:
        student_answer = answers.get(str(question.id), '')
        if student_answer.upper() == question.correct_answer.upper():
            correct += 1

    score = round((correct / total) * 100, 1) if total > 0 else 0

    submission = QuizSubmission.objects.create(
        quiz=quiz,
        student=request.user,
        answers=answers,
        score=score,
        total_questions=total,
        correct_count=correct,
    )

    # 返回结果和正确答案
    result_questions = []
    for q in questions:
        result_questions.append({
            'id': q.id,
            'question_text': q.question_text,
            'option_a': q.option_a,
            'option_b': q.option_b,
            'option_c': q.option_c,
            'option_d': q.option_d,
            'correct_answer': q.correct_answer,
            'explanation': q.explanation,
            'your_answer': answers.get(str(q.id), ''),
            'is_correct': answers.get(str(q.id), '').upper() == q.correct_answer.upper(),
        })

    return Response({
        'score': score,
        'total_questions': total,
        'correct_count': correct,
        'questions': result_questions,
        'attempt_number': attempts_used + 1,
        'remaining_attempts': max(0, quiz.max_attempts - attempts_used - 1),
    })


@api_view(['GET'])
@permission_classes([IsAuthenticated])
def quiz_submissions(request, quiz_id):
    """教师查看Quiz的所有提交"""
    try:
        quiz = Quiz.objects.get(id=quiz_id, creator=request.user)
    except Quiz.DoesNotExist:
        return Response({'error': 'Quiz不存在或无权操作'}, status=status.HTTP_404_NOT_FOUND)

    submissions = quiz.submissions.all().select_related('student')
    serializer = QuizSubmissionSerializer(submissions, many=True)
    return Response(serializer.data)


@api_view(['GET'])
@permission_classes([IsAuthenticated])
def quiz_statistics(request, quiz_id):
    """教师查看Quiz统计分析"""
    try:
        quiz = Quiz.objects.get(id=quiz_id, creator=request.user)
    except Quiz.DoesNotExist:
        return Response({'error': 'Quiz不存在或无权操作'}, status=status.HTTP_404_NOT_FOUND)

    submissions = quiz.submissions.all().select_related('student')
    questions = quiz.questions.all().order_by('order')

    # 每道题的统计
    question_stats = []
    for q in questions:
        total_answers = 0
        correct_answers = 0
        option_distribution = {'A': 0, 'B': 0, 'C': 0, 'D': 0}

        for sub in submissions:
            student_answer = sub.answers.get(str(q.id), '')
            if student_answer:
                total_answers += 1
                upper_answer = student_answer.upper()
                if upper_answer in option_distribution:
                    option_distribution[upper_answer] += 1
                if upper_answer == q.correct_answer.upper():
                    correct_answers += 1

        wrong_count = total_answers - correct_answers
        question_stats.append({
            'question_id': q.id,
            'question_text': q.question_text,
            'order': q.order,
            'correct_answer': q.correct_answer,
            'total_answers': total_answers,
            'correct_count': correct_answers,
            'wrong_count': wrong_count,
            'correct_rate': round((correct_answers / total_answers) * 100, 1) if total_answers > 0 else 0,
            'wrong_rate': round((wrong_count / total_answers) * 100, 1) if total_answers > 0 else 0,
            'option_distribution': option_distribution,
        })

    # 总体统计
    total_submissions = submissions.count()
    unique_students = submissions.values('student').distinct().count()
    avg_score = submissions.aggregate(avg=Avg('score'))['avg'] or 0
    score_distribution = {
        '0-59': submissions.filter(score__lt=60).count(),
        '60-69': submissions.filter(score__gte=60, score__lt=70).count(),
        '70-79': submissions.filter(score__gte=70, score__lt=80).count(),
        '80-89': submissions.filter(score__gte=80, score__lt=90).count(),
        '90-100': submissions.filter(score__gte=90).count(),
    }

    # 每个学生的所有提交记录
    student_submissions = []
    for sub in submissions:
        student_submissions.append({
            'student_name': sub.student.username,
            'score': sub.score,
            'correct_count': sub.correct_count,
            'total_questions': sub.total_questions,
            'submitted_at': sub.submitted_at.isoformat(),
            'answers': sub.answers,
        })

    return Response({
        'quiz_title': quiz.title,
        'max_attempts': quiz.max_attempts,
        'total_submissions': total_submissions,
        'unique_students': unique_students,
        'average_score': round(avg_score, 1),
        'score_distribution': score_distribution,
        'question_stats': question_stats,
        'student_submissions': student_submissions,
    })


@api_view(['GET'])
@permission_classes([IsAuthenticated])
def my_quiz_submissions(request, quiz_id):
    """学生查看自己在某个Quiz的所有提交记录"""
    try:
        quiz = Quiz.objects.get(id=quiz_id)
    except Quiz.DoesNotExist:
        return Response({'error': 'Quiz不存在'}, status=status.HTTP_404_NOT_FOUND)

    submissions = QuizSubmission.objects.filter(
        quiz=quiz, student=request.user
    ).order_by('-submitted_at')

    questions = quiz.questions.all().order_by('order')

    results = []
    for sub in submissions:
        question_details = []
        for q in questions:
            student_answer = sub.answers.get(str(q.id), '')
            question_details.append({
                'id': q.id,
                'question_text': q.question_text,
                'option_a': q.option_a,
                'option_b': q.option_b,
                'option_c': q.option_c,
                'option_d': q.option_d,
                'correct_answer': q.correct_answer,
                'explanation': q.explanation,
                'your_answer': student_answer,
                'is_correct': student_answer.upper() == q.correct_answer.upper() if student_answer else False,
            })
        results.append({
            'id': sub.id,
            'score': sub.score,
            'total_questions': sub.total_questions,
            'correct_count': sub.correct_count,
            'submitted_at': sub.submitted_at.isoformat(),
            'questions': question_details,
        })

    return Response({
        'quiz_title': quiz.title,
        'max_attempts': quiz.max_attempts,
        'attempts_used': submissions.count(),
        'remaining_attempts': max(0, quiz.max_attempts - submissions.count()),
        'submissions': results,
    })


@api_view(['DELETE', 'POST'])
@permission_classes([IsAuthenticated])
def delete_quiz(request, quiz_id):
    """删除Quiz"""
    try:
        quiz = Quiz.objects.get(id=quiz_id, creator=request.user)
    except Quiz.DoesNotExist:
        return Response({'error': 'Quiz不存在或无权操作'}, status=status.HTTP_404_NOT_FOUND)

    # 删除关联的源文件
    if quiz.source_file:
        try:
            quiz.source_file.delete(save=False)
        except Exception:
            pass  # 文件不存在也继续删除数据库记录
    
    quiz_title = quiz.title
    quiz.delete()
    return Response({'message': f'Quiz「{quiz_title}」已删除'}, status=status.HTTP_200_OK)


@api_view(['GET'])
@permission_classes([IsAuthenticated])
def pending_quizzes(request):
    """
    获取当前学生所有未完成的Quiz（用于日历DDL提醒）
    条件：
    1. Quiz已发布
    2. 截止时间未过（或无截止时间）
    3. 该学生尚未用完所有答题机会
    """
    from courses.models import Enrollment

    # 获取学生已选课程ID
    enrolled_course_ids = Enrollment.objects.filter(
        user=request.user
    ).values_list('course_id', flat=True)

    # 获取这些课程下的已发布Quiz（包含截止时间的）
    now = timezone.now()
    quizzes = Quiz.objects.filter(
        course_id__in=enrolled_course_ids,
        is_published=True,
        end_time__isnull=False,
        end_time__gt=now,
    ).select_related('course')

    # 过滤掉学生已用完答题机会的Quiz
    result = []
    for quiz in quizzes:
        attempts_used = QuizSubmission.objects.filter(
            quiz=quiz, student=request.user
        ).count()
        if attempts_used < quiz.max_attempts:
            result.append({
                'id': quiz.id,
                'title': quiz.title,
                'end_time': quiz.end_time.isoformat(),
                'course_id': quiz.course_id,
                'course_name': quiz.course.title if quiz.course else '',
                'share_code': quiz.share_code,
                'attempts_used': attempts_used,
                'max_attempts': quiz.max_attempts,
            })

    return Response(result)


@api_view(['GET'])
@permission_classes([IsAuthenticated])
def course_quizzes(request, course_id):
    """获取课程下的所有已发布Quiz（学生端）"""
    try:
        course = Course.objects.get(id=course_id)
    except Course.DoesNotExist:
        return Response({'error': '课程不存在'}, status=status.HTTP_404_NOT_FOUND)

    quizzes = Quiz.objects.filter(course=course, is_published=True)
    serializer = QuizStudentSerializer(quizzes, many=True, context={'request': request})
    return Response(serializer.data)


@api_view(['GET'])
@permission_classes([IsAuthenticated])
def teacher_analytics(request):
    """教师数据分析总览 - 获取该教师所有课程的Quiz答题情况"""
    courses = Course.objects.filter(instructor=request.user)
    if not courses.exists():
        return Response({'courses': [], 'overview': {
            'total_courses': 0, 'total_quizzes': 0,
            'total_submissions': 0, 'overall_avg_score': 0,
        }})

    all_quizzes = Quiz.objects.filter(course__in=courses).select_related('course')
    all_submissions = QuizSubmission.objects.filter(quiz__in=all_quizzes)

    # 总览
    overview = {
        'total_courses': courses.count(),
        'total_quizzes': all_quizzes.count(),
        'total_submissions': all_submissions.count(),
        'overall_avg_score': round(all_submissions.aggregate(avg=Avg('score'))['avg'] or 0, 1),
    }

    # 按课程分组
    course_data = []
    for course in courses:
        quizzes = all_quizzes.filter(course=course)
        course_submissions = all_submissions.filter(quiz__in=quizzes)

        quiz_list = []
        for quiz in quizzes:
            subs = course_submissions.filter(quiz=quiz)
            sub_count = subs.count()
            avg = round(subs.aggregate(a=Avg('score'))['a'] or 0, 1)

            score_dist = {
                '0-59': subs.filter(score__lt=60).count(),
                '60-69': subs.filter(score__gte=60, score__lt=70).count(),
                '70-79': subs.filter(score__gte=70, score__lt=80).count(),
                '80-89': subs.filter(score__gte=80, score__lt=90).count(),
                '90-100': subs.filter(score__gte=90).count(),
            }

            # 每题正确率
            questions = quiz.questions.all().order_by('order')
            question_stats = []
            for q in questions:
                total = 0
                correct = 0
                for sub in subs:
                    ans = sub.answers.get(str(q.id), '')
                    if ans:
                        total += 1
                        if ans.upper() == q.correct_answer.upper():
                            correct += 1
                question_stats.append({
                    'order': q.order,
                    'question_text': q.question_text[:50],
                    'correct_rate': round((correct / total) * 100, 1) if total > 0 else 0,
                    'total_answers': total,
                })

            # 学生提交列表
            student_records = []
            for sub in subs.select_related('student').order_by('-submitted_at'):
                student_records.append({
                    'student_name': sub.student.username,
                    'score': sub.score,
                    'correct_count': sub.correct_count,
                    'total_questions': sub.total_questions,
                    'submitted_at': sub.submitted_at.isoformat(),
                })

            quiz_list.append({
                'id': quiz.id,
                'title': quiz.title,
                'is_published': quiz.is_published,
                'submission_count': sub_count,
                'average_score': avg,
                'score_distribution': score_dist,
                'question_stats': question_stats,
                'student_records': student_records,
            })

        course_data.append({
            'id': course.id,
            'title': course.title,
            'quiz_count': quizzes.count(),
            'total_submissions': course_submissions.count(),
            'average_score': round(course_submissions.aggregate(a=Avg('score'))['a'] or 0, 1),
            'quizzes': quiz_list,
        })

    return Response({'overview': overview, 'courses': course_data})


@api_view(['POST'])
@permission_classes([IsAuthenticated])
def debate_start(request):
    """开启AI辩论对战"""
    serializer = DebateStartSerializer(data=request.data)
    if not serializer.is_valid():
        return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)

    topic = serializer.validated_data.get('topic', '').strip() or random.choice(DEBATE_TOPICS)
    course = None
    course_id = serializer.validated_data.get('course_id')
    if course_id:
        try:
            course = Course.objects.get(id=course_id)
        except Course.DoesNotExist:
            return Response({'error': '课程不存在'}, status=status.HTTP_404_NOT_FOUND)

    course_context = _build_course_context(course)
    ai_claim = _generate_debate_claim(topic, course_context)

    match = DebateMatch.objects.create(
        student=request.user,
        topic=topic,
        ai_claim=ai_claim,
        course=course,
    )

    return Response({
        'match_id': match.id,
        'topic': match.topic,
        'ai_claim': match.ai_claim,
        'status': match.status,
        'rounds_count': match.rounds_count,
        'total_attack': match.total_attack,
        'best_attack': match.best_attack,
        'course_title': course.title if course else '',
    }, status=status.HTTP_201_CREATED)


@api_view(['POST'])
@permission_classes([IsAuthenticated])
def debate_attack(request, match_id):
    """学生提交反驳，AI评分并反击"""
    try:
        match = DebateMatch.objects.get(id=match_id, student=request.user)
    except DebateMatch.DoesNotExist:
        return Response({'error': '对战不存在'}, status=status.HTTP_404_NOT_FOUND)

    if match.status != 'ongoing':
        return Response({'error': '该对战已结束，请开启新挑战'}, status=status.HTTP_400_BAD_REQUEST)

    serializer = DebateAttackSerializer(data=request.data)
    if not serializer.is_valid():
        return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)

    argument = serializer.validated_data['argument'].strip()
    if len(argument) < 20:
        return Response({'error': '反驳内容太短，至少输入20个字'}, status=status.HTTP_400_BAD_REQUEST)

    round_number = match.rounds.count() + 1
    course_context = _build_course_context(match.course)
    score_data = _evaluate_argument(argument, match.topic, match.ai_claim, course_context)
    ai_counter = _generate_ai_counter(match.topic, match.ai_claim, argument, score_data)

    debate_round = DebateRound.objects.create(
        match=match,
        round_number=round_number,
        student_argument=argument,
        ai_counter=ai_counter,
        logic_score=score_data['logic_score'],
        evidence_score=score_data['evidence_score'],
        knowledge_score=score_data['knowledge_score'],
        structure_score=score_data['structure_score'],
        attack_power=score_data['attack_power'],
        verdict=score_data['verdict'],
    )

    match.rounds_count = round_number
    match.total_attack += debate_round.attack_power
    match.best_attack = max(match.best_attack, debate_round.attack_power)

    avg_attack = match.total_attack / max(1, match.rounds_count)
    if debate_round.attack_power >= 88 or match.total_attack >= 230 or (round_number >= 3 and avg_attack >= 78):
        match.status = 'won'
    elif round_number >= 5 and avg_attack < 68:
        match.status = 'lost'

    match.save()

    gained_badges = []
    if match.status == 'won':
        gained_badges = _grant_debate_badges(request.user, match, debate_round)

    return Response({
        'match_id': match.id,
        'round_number': round_number,
        'attack_power': debate_round.attack_power,
        'score_breakdown': {
            'logic': debate_round.logic_score,
            'evidence': debate_round.evidence_score,
            'knowledge': debate_round.knowledge_score,
            'structure': debate_round.structure_score,
        },
        'verdict': debate_round.verdict,
        'ai_counter': debate_round.ai_counter,
        'match_status': match.status,
        'total_attack': match.total_attack,
        'best_attack': match.best_attack,
        'gained_badges': DebateBadgeSerializer(gained_badges, many=True).data,
    })


@api_view(['POST'])
@permission_classes([IsAuthenticated])
def debate_quit(request, match_id):
    """退出当前辩论对战"""
    try:
        match = DebateMatch.objects.get(id=match_id, student=request.user)
    except DebateMatch.DoesNotExist:
        return Response({'error': '对战不存在'}, status=status.HTTP_404_NOT_FOUND)

    if match.status != 'ongoing':
        return Response({'error': '该对战已结束'}, status=status.HTTP_400_BAD_REQUEST)

    match.status = 'lost'
    match.save(update_fields=['status', 'updated_at'])
    return Response({'message': '已退出当前对战', 'match_id': match.id, 'status': match.status})


@api_view(['DELETE'])
@permission_classes([IsAuthenticated])
def debate_delete(request, match_id):
    """隐藏指定辩论对战记录（不影响累计统计）"""
    try:
        match = DebateMatch.objects.get(id=match_id, student=request.user)
    except DebateMatch.DoesNotExist:
        return Response({'error': '对战不存在'}, status=status.HTTP_404_NOT_FOUND)

    if match.is_hidden:
        return Response({'message': '战绩已隐藏', 'match_id': match_id}, status=status.HTTP_200_OK)

    match.is_hidden = True
    match.save(update_fields=['is_hidden', 'updated_at'])
    return Response({'message': '战绩已隐藏', 'match_id': match_id}, status=status.HTTP_200_OK)


@api_view(['GET'])
@permission_classes([IsAuthenticated])
def debate_profile(request):
    """获取当前用户辩论场战绩与勋章"""
    matches = DebateMatch.objects.filter(student=request.user)
    total_matches = matches.count()
    wins = matches.filter(status='won').count()
    best_attack = matches.aggregate(v=Avg('best_attack'))['v'] or 0

    visible_matches = matches.filter(is_hidden=False)
    latest_ongoing = visible_matches.filter(status='ongoing').first()
    recent_matches = visible_matches.order_by('-created_at')[:8]
    recent_data = [
        {
            'id': m.id,
            'topic': m.topic,
            'status': m.status,
            'best_attack': m.best_attack,
            'rounds_count': m.rounds_count,
            'created_at': m.created_at.isoformat(),
        }
        for m in recent_matches
    ]

    badges = DebateBadge.objects.filter(student=request.user)
    return Response({
        'summary': {
            'total_matches': total_matches,
            'wins': wins,
            'win_rate': round((wins / total_matches) * 100, 1) if total_matches else 0,
            'average_best_attack': round(best_attack, 1),
            'badge_count': badges.count(),
        },
        'current_match': {
            'id': latest_ongoing.id,
            'topic': latest_ongoing.topic,
            'ai_claim': latest_ongoing.ai_claim,
            'rounds_count': latest_ongoing.rounds_count,
            'status': latest_ongoing.status,
            'total_attack': latest_ongoing.total_attack,
            'best_attack': latest_ongoing.best_attack,
            'rounds': [
                {
                    'id': r.id,
                    'round_number': r.round_number,
                    'student_argument': r.student_argument,
                    'ai_counter': r.ai_counter,
                    'attack_power': r.attack_power,
                    'verdict': r.verdict,
                }
                for r in latest_ongoing.rounds.all().order_by('round_number')
            ],
        } if latest_ongoing else None,
        'recent_matches': recent_data,
        'badges': DebateBadgeSerializer(badges, many=True).data,
    })


# ---------- Quiz 邮件提醒 API ----------

@api_view(['GET'])
@permission_classes([IsAuthenticated])
def quiz_reminder_logs(request, quiz_id):
    """GET /api/ai/quiz/<quiz_id>/reminder-logs/ — 查看该 Quiz 的提醒发送记录"""
    from .models import QuizReminderLog

    try:
        quiz = Quiz.objects.get(id=quiz_id, creator=request.user)
    except Quiz.DoesNotExist:
        return Response({'error': 'Quiz不存在或无权操作'}, status=404)

    logs = QuizReminderLog.objects.filter(quiz=quiz).select_related('student').order_by('-sent_at')
    data = [
        {
            'student_id': log.student.id,
            'student_name': log.student.username,
            'student_email': log.student.email,
            'sent_at': log.sent_at.isoformat(),
        }
        for log in logs
    ]
    return Response({'quiz_id': quiz_id, 'quiz_title': quiz.title, 'reminder_logs': data})


@api_view(['POST'])
@permission_classes([IsAuthenticated])
def send_quiz_reminder_now(request, quiz_id):
    """POST /api/ai/quiz/<quiz_id>/send-reminders/ — 立即向未完成学生发送提醒邮件"""
    from django.core.mail import send_mail
    from .models import QuizReminderLog, QuizSubmission
    from courses.models import Enrollment

    try:
        quiz = Quiz.objects.get(id=quiz_id, creator=request.user)
    except Quiz.DoesNotExist:
        return Response({'error': 'Quiz不存在或无权操作'}, status=404)

    if not quiz.is_published:
        return Response({'error': 'Quiz尚未发布，无法发送提醒'}, status=400)

    if quiz.end_time and timezone.now() > quiz.end_time:
        return Response({'error': 'Quiz已截止，无法发送提醒'}, status=400)

    if not quiz.course:
        return Response({'error': 'Quiz未关联课程，无法确定学生名单'}, status=400)

    platform_email = settings.EMAIL_HOST_USER
    if not platform_email:
        return Response({'error': '平台邮箱未配置（EMAIL_HOST_USER）'}, status=503)

    frontend_url = getattr(settings, 'FRONTEND_URL', 'http://localhost:3000')
    enrollments = Enrollment.objects.filter(course=quiz.course).select_related('user')
    submitted_ids = set(QuizSubmission.objects.filter(quiz=quiz).values_list('student_id', flat=True))
    reminded_ids = set(QuizReminderLog.objects.filter(quiz=quiz).values_list('student_id', flat=True))

    sent, skipped, failed = 0, 0, 0

    for enrollment in enrollments:
        student = enrollment.user
        if student.id in submitted_ids or student.id in reminded_ids:
            skipped += 1
            continue
        if not student.email:
            skipped += 1
            continue

        end_time_str = ''
        if quiz.end_time:
            end_time_local = timezone.localtime(quiz.end_time)
            end_time_str = end_time_local.strftime('%Y年%m月%d日 %H:%M')

        subject = f'【提醒】Quiz「{quiz.title}」' + (f'将于 {end_time_str} 截止' if end_time_str else '尚未完成')
        body = (
            f'{student.username} 同学，您好！\n\n'
            f'您在课程「{quiz.course.title}」中有一份 Quiz 尚未完成，请尽快完成：\n\n'
            f'  📝 Quiz 名称：{quiz.title}\n'
            + (f'  ⏰ 截止时间：{end_time_str}\n' if end_time_str else '')
            + f'  🔗 答题链接：{frontend_url}/quiz/{quiz.share_code}\n\n'
            f'请尽快完成，祝学习顺利！\n\n'
            f'—— {quiz.creator.username} 老师'
        )

        try:
            send_mail(
                subject=subject,
                message=body,
                from_email=platform_email,
                recipient_list=[student.email],
                fail_silently=False,
            )
            QuizReminderLog.objects.get_or_create(quiz=quiz, student=student)
            sent += 1
        except Exception as e:
            failed += 1
            logger.error(f'[Quiz提醒] 手动发送失败 → {student.email}，原因: {e}')

    return Response({
        'quiz_id': quiz_id,
        'quiz_title': quiz.title,
        'sent': sent,
        'skipped': skipped,
        'failed': failed,
    })


# ============ RAG 课程知识库 API ============

@api_view(['POST'])
@permission_classes([IsAuthenticated])
def rag_ask(request):
    """学生基于课程材料提问"""
    course_id = request.data.get('course_id')
    question = (request.data.get('question') or '').strip()

    if not course_id:
        return Response({'error': 'course_id 必填', 'code': 'MISSING_COURSE_ID'}, status=400)
    if not question:
        return Response({'error': '问题不能为空', 'code': 'EMPTY_QUESTION'}, status=400)

    try:
        course_id = int(course_id)
    except (TypeError, ValueError):
        return Response({'error': 'course_id 必须为整数', 'code': 'INVALID_COURSE_ID'}, status=400)

    if not Course.objects.filter(id=course_id).exists():
        return Response({'error': '课程不存在', 'code': 'COURSE_NOT_FOUND'}, status=404)

    from .rag import ask_course
    result = ask_course(course_id, question)
    return Response(result)


@api_view(['POST'])
@permission_classes([IsAuthenticated])
def rag_build_index(request):
    """教师重建某门课的 RAG 索引（仅讲师本人或超管可调用）"""
    course_id = request.data.get('course_id')
    if not course_id:
        return Response({'error': 'course_id 必填', 'code': 'MISSING_COURSE_ID'}, status=400)
    try:
        course_id = int(course_id)
    except (TypeError, ValueError):
        return Response({'error': 'course_id 必须为整数', 'code': 'INVALID_COURSE_ID'}, status=400)

    course = Course.objects.filter(id=course_id).first()
    if not course:
        return Response({'error': '课程不存在', 'code': 'COURSE_NOT_FOUND'}, status=404)

    if not (request.user.is_superuser or course.instructor_id == request.user.id):
        return Response({'error': '只有该课程的讲师可以重建索引', 'code': 'FORBIDDEN'}, status=403)

    from .rag import build_course_index
    try:
        result = build_course_index(course_id)
    except Exception as e:
        logger.exception('build_course_index failed')
        return Response({'error': f'建索引失败: {e}', 'code': 'BUILD_ERROR'}, status=500)
    return Response(result)


# ---------- 学习助手 Agent ----------

@api_view(['POST'])
@permission_classes([IsAuthenticated])
def agent_run(request):
    """
    POST /api/ai/agent/run/
    Body: { "message": "...", "history": [{"role":"user|assistant","content":"..."}] }
    返回: { "output", "steps", "elapsed_sec" } 或 { "error", "code", "elapsed_sec" }
    """
    message = request.data.get('message')
    history = request.data.get('history') or []
    if not isinstance(message, str) or not message.strip():
        return Response({'error': 'message 必填', 'code': 'EMPTY_MESSAGE', 'elapsed_sec': 0.0})
    if not isinstance(history, list):
        history = []

    from .agents import run_student_agent
    result = run_student_agent(request.user, message, history)
    return Response(result)

