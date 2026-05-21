"""
课程 RAG 知识库
功能：把某门课的所有 CourseFile 解析成文本 → 切片 → 向量化 → 存到 FAISS
每门课一个独立索引：backend/vector_db/course_<id>/
"""
import logging
from pathlib import Path
from typing import List, Dict, Optional

from django.conf import settings

logger = logging.getLogger(__name__)

# 全局单例：embedding 模型 + LangChain 组件（懒加载）
_embeddings_singleton = None


# ========== 文档解析 ==========

def _load_pdf(path: Path) -> List[Dict]:
    """返回 [{'text': str, 'page': int}, ...]"""
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    pages = []
    for i, page in enumerate(reader.pages):
        try:
            text = page.extract_text() or ''
        except Exception as e:
            logger.warning(f"PDF 第 {i+1} 页解析失败: {e}")
            text = ''
        if text.strip():
            pages.append({'text': text, 'page': i + 1})
    return pages


def _load_docx(path: Path) -> List[Dict]:
    """Word .docx：合并所有段落文字"""
    from docx import Document
    doc = Document(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    # 也提取表格里的文字
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    paragraphs.append(cell.text)
    text = '\n'.join(paragraphs)
    return [{'text': text, 'page': 1}] if text.strip() else []


def _load_pptx(path: Path) -> List[Dict]:
    """PowerPoint .pptx：每张幻灯片一个段，含文本框 + 备注"""
    from pptx import Presentation
    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = ''.join(run.text for run in para.runs)
                    if txt.strip():
                        parts.append(txt)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            parts.append(cell.text)
        # 演讲者备注
        if slide.has_notes_slide:
            note = slide.notes_slide.notes_text_frame.text
            if note.strip():
                parts.append(f'[备注] {note}')
        text = '\n'.join(parts)
        if text.strip():
            slides.append({'text': text, 'page': i + 1})
    return slides


def _load_text(path: Path) -> List[Dict]:
    """.md 或 .txt 直接读取"""
    try:
        text = path.read_text(encoding='utf-8', errors='ignore')
    except Exception as e:
        logger.warning(f"文本文件读取失败: {e}")
        return []
    if path.suffix.lower() == '.md':
        # 简单去除 markdown 标记
        import re
        text = re.sub(r'```[\s\S]*?```', '', text)  # 代码块
        text = re.sub(r'[#*_`>]', '', text)
    return [{'text': text, 'page': 1}] if text.strip() else []


_LOADERS = {
    '.pdf': _load_pdf,
    '.docx': _load_docx,
    '.pptx': _load_pptx,
    '.md': _load_text,
    '.txt': _load_text,
}


def _load_document(file_path: Path) -> List[Dict]:
    """根据后缀分发；返回 [{'text', 'page'}, ...]，失败返回 []"""
    ext = file_path.suffix.lower()
    loader = _LOADERS.get(ext)
    if not loader:
        logger.info(f"跳过不支持的文件: {file_path.name}")
        return []
    try:
        return loader(file_path)
    except Exception as e:
        logger.warning(f"解析失败 {file_path.name}: {e}")
        return []


# ========== 切片 + 向量化 ==========

def get_embeddings():
    """单例：本地 HuggingFace embedding 模型"""
    global _embeddings_singleton
    if _embeddings_singleton is None:
        from langchain_community.embeddings import HuggingFaceEmbeddings
        logger.info("加载本地 embedding 模型: paraphrase-multilingual-MiniLM-L12-v2")
        _embeddings_singleton = HuggingFaceEmbeddings(
            model_name='paraphrase-multilingual-MiniLM-L12-v2'
        )
    return _embeddings_singleton


def _split_documents(raw_docs: List[Dict], course_id: int, source_name: str):
    """切片，返回 LangChain Document 列表"""
    from langchain_core.documents import Document
    from langchain_text_splitters import RecursiveCharacterTextSplitter

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=settings.RAG_CHUNK_SIZE,
        chunk_overlap=settings.RAG_CHUNK_OVERLAP,
        separators=['\n\n', '\n', '。', '.', ' ', ''],
    )
    out = []
    for d in raw_docs:
        for chunk in splitter.split_text(d['text']):
            out.append(Document(
                page_content=chunk,
                metadata={
                    'course_id': course_id,
                    'source': source_name,
                    'page': d.get('page', 1),
                },
            ))
    return out


def _course_index_dir(course_id: int) -> Path:
    return Path(settings.RAG_VECTOR_DB_ROOT) / f'course_{course_id}'


# ========== 主入口 ==========

def build_course_index(course_id: int) -> Dict:
    """
    全量重建某门课的 FAISS 索引。
    返回统计：{'files': N, 'chunks': N, 'skipped': [...]}
    """
    from langchain_community.vectorstores import FAISS
    from courses.models import CourseFile

    files = CourseFile.objects.filter(course_id=course_id, file__isnull=False)
    if not files.exists():
        return {'files': 0, 'chunks': 0, 'skipped': [], 'message': '该课程无文件'}

    all_chunks = []
    skipped = []
    processed = 0

    for cf in files:
        if not cf.file or not cf.file.path:
            continue
        path = Path(cf.file.path)
        if not path.exists():
            skipped.append(f'{cf.file_name}（文件不存在）')
            continue
        if path.suffix.lower() not in settings.RAG_SUPPORTED_EXTS:
            skipped.append(f'{cf.file_name}（不支持的格式）')
            continue

        raw_docs = _load_document(path)
        if not raw_docs:
            skipped.append(f'{cf.file_name}（无法提取文字）')
            continue

        # 截断超大文档
        total = sum(len(d['text']) for d in raw_docs)
        if total > settings.RAG_MAX_DOC_CHARS:
            logger.warning(f'{cf.file_name} 超过 {settings.RAG_MAX_DOC_CHARS} 字，将截断')
            kept, used = [], 0
            for d in raw_docs:
                if used + len(d['text']) > settings.RAG_MAX_DOC_CHARS:
                    d = {**d, 'text': d['text'][: settings.RAG_MAX_DOC_CHARS - used]}
                    kept.append(d)
                    break
                kept.append(d)
                used += len(d['text'])
            raw_docs = kept

        chunks = _split_documents(raw_docs, course_id, cf.file_name)
        all_chunks.extend(chunks)
        processed += 1
        logger.info(f'{cf.file_name}: {len(chunks)} 个片段')

    if not all_chunks:
        return {'files': processed, 'chunks': 0, 'skipped': skipped, 'message': '没有可索引的内容'}

    embeddings = get_embeddings()
    vector_store = FAISS.from_documents(all_chunks, embeddings)

    out_dir = _course_index_dir(course_id)
    out_dir.mkdir(parents=True, exist_ok=True)
    vector_store.save_local(str(out_dir))
    logger.info(f'课程 {course_id} 索引已保存到 {out_dir}')

    # 更新索引后清掉缓存，下次问答会重新加载
    _vector_store_cache.pop(course_id, None)

    return {
        'files': processed,
        'chunks': len(all_chunks),
        'skipped': skipped,
        'index_dir': str(out_dir),
    }


def add_file_to_course_index(course_id: int, file_path, source_name: str = None) -> Dict:
    """
    增量：把单个文件追加进某门课的 FAISS 索引。
    - 索引已存在：load → add_documents → save
    - 索引不存在：直接新建只含本文件的索引
    返回 {'chunks': N, 'source': name, 'mode': 'append'|'create'} 或带 'skipped' 的失败信息
    """
    from langchain_community.vectorstores import FAISS

    path = Path(file_path)
    if not path.exists():
        return {'chunks': 0, 'skipped': f'{path.name}（文件不存在）'}
    if path.suffix.lower() not in settings.RAG_SUPPORTED_EXTS:
        return {'chunks': 0, 'skipped': f'{path.name}（不支持的格式）'}

    raw_docs = _load_document(path)
    if not raw_docs:
        return {'chunks': 0, 'skipped': f'{path.name}（无法提取文字）'}

    # 截断超大文档
    total = sum(len(d['text']) for d in raw_docs)
    if total > settings.RAG_MAX_DOC_CHARS:
        logger.warning(f'{path.name} 超过 {settings.RAG_MAX_DOC_CHARS} 字，将截断')
        kept, used = [], 0
        for d in raw_docs:
            if used + len(d['text']) > settings.RAG_MAX_DOC_CHARS:
                d = {**d, 'text': d['text'][: settings.RAG_MAX_DOC_CHARS - used]}
                kept.append(d)
                break
            kept.append(d)
            used += len(d['text'])
        raw_docs = kept

    name = source_name or path.name
    chunks = _split_documents(raw_docs, course_id, name)
    if not chunks:
        return {'chunks': 0, 'skipped': f'{name}（切片为空）'}

    embeddings = get_embeddings()
    out_dir = _course_index_dir(course_id)
    index_file = out_dir / 'index.faiss'

    if index_file.exists():
        vs = FAISS.load_local(
            str(out_dir), embeddings, allow_dangerous_deserialization=True,
        )
        vs.add_documents(chunks)
        mode = 'append'
    else:
        out_dir.mkdir(parents=True, exist_ok=True)
        vs = FAISS.from_documents(chunks, embeddings)
        mode = 'create'

    vs.save_local(str(out_dir))
    _vector_store_cache.pop(course_id, None)
    logger.info(f'[RAG] {mode} course={course_id} source={name} chunks={len(chunks)}')
    return {'chunks': len(chunks), 'source': name, 'mode': mode, 'index_dir': str(out_dir)}


# ========== 问答（RAG 检索 + LLM 生成） ==========

# 课程 vector_store 缓存：{course_id: FAISS}
_vector_store_cache: Dict[int, object] = {}


def _load_course_index(course_id: int):
    """加载某门课的 FAISS 索引（带缓存）。索引不存在返回 None。"""
    if course_id in _vector_store_cache:
        return _vector_store_cache[course_id]

    index_dir = _course_index_dir(course_id)
    if not (index_dir / 'index.faiss').exists():
        return None

    from langchain_community.vectorstores import FAISS
    embeddings = get_embeddings()
    vs = FAISS.load_local(
        str(index_dir),
        embeddings,
        allow_dangerous_deserialization=True,  # 我们自己生成的 pkl，可信
    )
    _vector_store_cache[course_id] = vs
    return vs


def _build_prompt(question: str, hits: list) -> str:
    """把检索结果拼成 prompt"""
    blocks = []
    for i, doc in enumerate(hits, 1):
        meta = doc.metadata or {}
        src = meta.get('source', '?')
        page = meta.get('page', '?')
        blocks.append(f'[材料{i}]（来源: {src} 第 {page} 页）\n{doc.page_content}')

    materials = '\n\n'.join(blocks) if blocks else '（无相关材料）'
    return (
        '你是课程助教。请仅根据下面提供的"课程材料"回答学生的问题。\n'
        '规则：\n'
        '1. 如果材料里没有答案，直接回答"根据现有课程材料无法回答"，不要编造。\n'
        '2. 引用具体内容时用 [材料N] 标注。\n'
        '3. 回答用学生问题相同的语言（中文问就中文答，英文问就英文答）。\n\n'
        f'【课程材料】\n{materials}\n\n'
        f'【学生问题】\n{question}\n\n'
        '【你的回答】\n'
    )


def ask_course(course_id: int, question: str, top_k: Optional[int] = None) -> Dict:
    """
    RAG 问答主入口。
    返回:
      成功: {answer, sources, course_id, question, model}
      失败: {error, code, course_id, question}
    """
    question = (question or '').strip()
    if not question:
        return {'error': '问题不能为空', 'code': 'EMPTY_QUESTION', 'course_id': course_id, 'question': question}

    vs = _load_course_index(course_id)
    if vs is None:
        return {
            'error': '该课程尚未建立知识库，请先让教师上传材料并建立索引',
            'code': 'NO_INDEX',
            'course_id': course_id,
            'question': question,
        }

    k = top_k or settings.RAG_TOP_K
    try:
        hits = vs.similarity_search(question, k=k)
    except Exception as e:
        logger.error(f'检索失败: {e}')
        return {'error': f'检索失败: {e}', 'code': 'SEARCH_ERROR', 'course_id': course_id, 'question': question}

    prompt = _build_prompt(question, hits)

    # 调用 LLM（复用项目现有 GitHub Models / OpenAI 配置）
    try:
        from langchain_openai import ChatOpenAI
        llm_kwargs = {
            'model': settings.AI_MODEL_NAME,
            'temperature': 0.3,
            'api_key': settings.OPENAI_API_KEY,
        }
        if getattr(settings, 'USE_GITHUB_MODELS', False) and getattr(settings, 'OPENAI_API_BASE', None):
            llm_kwargs['base_url'] = settings.OPENAI_API_BASE
        llm = ChatOpenAI(**llm_kwargs)
        from langchain_core.messages import HumanMessage
        resp = llm.invoke([HumanMessage(content=prompt)])
        answer = resp.content if hasattr(resp, 'content') else str(resp)
    except Exception as e:
        logger.error(f'LLM 调用失败: {e}')
        return {'error': f'AI 服务暂时不可用: {e}', 'code': 'LLM_ERROR', 'course_id': course_id, 'question': question}

    sources = []
    for doc in hits:
        meta = doc.metadata or {}
        snippet = (doc.page_content or '')[:120].replace('\n', ' ')
        sources.append({
            'file': meta.get('source', '?'),
            'page': meta.get('page', '?'),
            'snippet': snippet,
        })

    return {
        'answer': answer,
        'sources': sources,
        'course_id': course_id,
        'question': question,
        'model': settings.AI_MODEL_NAME,
    }
